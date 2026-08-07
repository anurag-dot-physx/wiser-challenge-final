from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

HERE = Path(__file__).resolve().parent
PPTX = HERE / "WISER_Vanguard_Quantum_Portfolio_Challenge_2026.pptx"

BG = RGBColor(247, 247, 248)
INK = RGBColor(32, 33, 38)
MUT = RGBColor(107, 107, 114)
BURG = RGBColor(123, 30, 53)
PUR = RGBColor(91, 42, 134)
ORG = RGBColor(227, 106, 54)
GREEN = RGBColor(43, 122, 75)
BLUE = RGBColor(49, 90, 138)
LBLUE = RGBColor(120, 184, 232)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(217, 217, 222)


def add_text(slide, text, x, y, w, h, size=16, bold=False, color=INK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    p.alignment = align
    return box


def rect(slide, x, y, w, h, fill, line=LINE, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    sh = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = line
    return sh


def blank_slide(prs):
    layout = next((lay for lay in prs.slide_layouts if "blank" in (lay.name or "").lower()), prs.slide_layouts[0])
    slide = prs.slides.add_slide(layout)
    for shape in list(slide.shapes):
        if getattr(shape, "is_placeholder", False):
            sp = shape._element
            sp.getparent().remove(sp)
    return slide


def base(slide, title, subtitle):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    rect(slide, 0, 0, 13.333, 0.08, BURG, BURG)
    add_text(slide, title, 0.6, 0.35, 12, 0.5, 26, True)
    add_text(slide, subtitle, 0.62, 0.92, 11.9, 0.35, 12, False, MUT)
    add_text(slide, "WISER x Vanguard | Quantum for Finance 2026", 0.55, 7.12, 5.5, 0.2, 8, False, MUT)


def bar(slide, x, baseline, width, height, color, label, value_text):
    rect(slide, x, baseline-height, width, height, color, color)
    add_text(slide, value_text, x-0.08, baseline-height-0.35, width+0.16, 0.25, 11, True, color, PP_ALIGN.CENTER)
    if label:
        add_text(slide, label, x-0.25, baseline+0.08, width+0.5, 0.35, 10.5, False, INK, PP_ALIGN.CENTER)


def add_flagship_slide(prs):
    s = blank_slide(prs)
    base(s, "Result 1: exact classical audit and turnover control", "The production model is validated exhaustively before any quantum benchmark is interpreted.")
    add_text(s, "QP → exact-grid objective gap", 0.7, 1.45, 4.7, 0.3, 18, True)
    rect(s, 0.8, 4.4, 5.0, 0.015, LINE, LINE)
    vals = [("Growth", 0.00081799, BLUE, "8.18e−4"), ("Balanced", 0.00066947, BURG, "6.69e−4"), ("Defensive", 0.00140581, PUR, "1.41e−3")]
    scale = 2.2 / 0.0015
    for i, (lab, v, c, txt) in enumerate(vals):
        bar(s, 1.15+i*1.45, 4.4, 0.7, v*scale, c, lab, txt)
    add_text(s, "Small positive gaps are expected from discretizing a convex relaxation.", 0.9, 5.05, 4.8, 0.55, 11, False, MUT, PP_ALIGN.CENTER)

    add_text(s, "One-way turnover", 6.7, 1.45, 3.8, 0.3, 18, True)
    rect(s, 6.85, 4.4, 5.3, 0.015, LINE, LINE)
    bar(s, 7.55, 4.4, 1.0, 2.25, BURG, "λT = 0", "40%")
    bar(s, 9.7, 4.4, 1.0, 0.84, GREEN, "λT = 2", "15%")
    add_text(s, "62.5% reduction", 9.0, 5.05, 2.7, 0.38, 18, True, GREEN, PP_ALIGN.CENTER)
    rect(s, 0.8, 5.75, 11.55, 0.72, RGBColor(238, 246, 241), RGBColor(183, 215, 196), True)
    add_text(s, "19,448 states enumerated • zero canonical hard breaches • reduced QUBO ground audited against exact reference", 1.05, 5.97, 11.05, 0.3, 13, True, GREEN, PP_ALIGN.CENTER)


def add_qaoa_slide(prs):
    s = blank_slide(prs)
    base(s, "Result 2: reduced QUBO/QAOA benchmark recovers the exact allocation", "Balanced reduced model: 15 total variables = 9 allocation + 6 slack; 66 hard-feasible portfolios.")
    labels = ["US Equity", "Intl Equity", "Govt Bonds", "Commodities", "Cash"]
    exact = [37.5, 12.5, 37.5, 12.5, 0]
    recovered = [37.5, 12.5, 37.5, 12.5, 0]
    baseline = 5.5
    rect(s, 0.75, baseline, 9.5, 0.015, LINE, LINE)
    for i, lab in enumerate(labels):
        x = 1.0+i*1.85
        if exact[i] > 0:
            bar(s, x, baseline, 0.42, exact[i]*0.065, LBLUE, "", f"{exact[i]:g}%")
            bar(s, x+0.45, baseline, 0.42, recovered[i]*0.065, BLUE, "", f"{recovered[i]:g}%")
        add_text(s, lab, x-0.35, 5.65, 1.65, 0.45, 10.5, False, INK, PP_ALIGN.CENTER)
    rect(s, 10.55, 1.75, 1.95, 1.42, WHITE, LINE, True)
    rect(s, 10.78, 2.08, 0.22, 0.22, LBLUE, LBLUE)
    add_text(s, "Reduced exact", 11.1, 2.02, 1.1, 0.28, 10.5)
    rect(s, 10.78, 2.47, 0.22, 0.22, BLUE, BLUE)
    add_text(s, "QAOA / fallback", 11.1, 2.41, 1.2, 0.28, 10.5)
    rect(s, 10.55, 3.55, 1.95, 1.35, RGBColor(238, 246, 241), RGBColor(183, 215, 196), True)
    add_text(s, "QUBO ground\naudited", 10.82, 3.82, 1.4, 0.5, 13, True, GREEN, PP_ALIGN.CENTER)
    add_text(s, "TRUE", 10.85, 4.35, 1.35, 0.35, 20, True, GREEN, PP_ALIGN.CENTER)
    rect(s, 0.8, 6.3, 11.55, 0.62, RGBColor(238, 246, 241), RGBColor(183, 215, 196), True)
    add_text(s, "The recovered allocation matches the reduced exact optimum; this benchmark is intentionally separated from the eight-asset production claim.", 1.05, 6.48, 11.0, 0.26, 12.5, True, GREEN, PP_ALIGN.CENTER)


def add_vqe_slide(prs):
    s = blank_slide(prs)
    base(s, "Result 3: higher-moment HUBO/VQE remains feasible but shows solver gap", "Exploratory 15-qubit benchmark. Co-skewness and co-kurtosis enrich the landscape; all shown states have zero budget breach.")
    names = ["Financial GT", "Exact HUBO", "H ground", "Finite-shot VQE"]
    cols = [BLUE, GREEN, PUR, ORG]
    returns = [53.41, 21.08, 21.08, 19.75]
    sharpes = [1.09482, 0.73544, 0.73544, 0.68444]
    add_text(s, "Expected return", 0.7, 1.45, 4.6, 0.3, 18, True)
    rect(s, 0.8, 4.55, 5.2, 0.015, LINE, LINE)
    for i, (n, v, c) in enumerate(zip(names, returns, cols)):
        bar(s, 1.0+i*1.22, 4.55, 0.62, v*0.048, c, n, f"{v:.2f}%")
    add_text(s, "Sharpe ratio", 6.45, 1.45, 4.6, 0.3, 18, True)
    rect(s, 6.55, 4.55, 5.9, 0.015, LINE, LINE)
    for i, (n, v, c) in enumerate(zip(names, sharpes, cols)):
        bar(s, 6.75+i*1.38, 4.55, 0.62, v*2.25, c, n, f"{v:.3f}")
    rect(s, 0.8, 5.85, 11.55, 0.82, RGBColor(255, 245, 239), RGBColor(239, 195, 173), True)
    add_text(s, "Finite-shot VQE is feasible but does not recover the exact HUBO ground state in this run. Reported as solver-quality evidence, not quantum advantage.", 1.05, 6.08, 11.05, 0.35, 13, True, BURG, PP_ALIGN.CENTER)


def main():
    prs = Presentation(PPTX)
    add_flagship_slide(prs)
    add_qaoa_slide(prs)
    add_vqe_slide(prs)
    prs.save(PPTX)
    print(f"Appended 3 result slides to {PPTX}")


if __name__ == "__main__":
    main()
